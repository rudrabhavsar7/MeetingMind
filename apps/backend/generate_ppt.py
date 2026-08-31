from pptx import Presentation

# Initialize presentation
prs = Presentation()


# Helper function to add title slide
def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text


# Helper function to add bullet slide
def add_bullet_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title.text = title_text

    tf = body_shape.text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0


# Slide 1: Title
add_title_slide(
    "MeetingMind", "Privacy-First, Self-Hosted AI Meeting Intelligence Platform\nInternal Review 1\n\nTeam: Rudra, Arnish, & Team"
)

# Slide 2: Problem Statement
add_bullet_slide(
    "Problem Statement",
    [
        "Cloud-based meeting bots pose massive privacy and security risks for enterprises.",
        "Internal meetings often contain highly sensitive financial, strategic, and personal data.",
        "Existing solutions force organizations to send their data to third-party servers.",
        "Organizations need the power of AI summaries and transcripts without compromising data sovereignty.",
    ],
)

# Slide 3: Proposed Solution: MeetingMind
add_bullet_slide(
    "Our Solution: MeetingMind",
    [
        "A 100% self-hosted meeting intelligence platform.",
        "Privacy-First: No audio or transcripts leave the operator-controlled infrastructure by default.",
        "Extension-First: A Chrome extension captures live meeting audio directly from the browser (e.g., Google Meet).",
        "Local AI Models: Utilizes local AI for transcription (Whisper), diarization, and LLM-based summarization (Ollama).",
    ],
)

# Slide 4: Key Features & Capabilities
add_bullet_slide(
    "Key Features",
    [
        "Real-Time Capture: Live PCM 16-bit LE audio ingestion via WebSockets.",
        "Speaker-Aware Transcripts: Accurate speaker diarization (who said what).",
        "AI Summarization: Executive summaries, action items, and key decisions.",
        "Search & RAG: Searchable transcript embeddings with exact citations across meeting history.",
        "Workspace Isolation: Strict data isolation for multi-tenant self-hosting.",
    ],
)

# Slide 5: System Architecture & Tech Stack
add_bullet_slide(
    "System Architecture & Tech Stack",
    [
        "Frontend: Next.js 15 App Router, React 19, Tailwind CSS v4, Zustand.",
        "Capture: Chrome Extension (Manifest V3) using Offscreen API.",
        "Backend API: FastAPI, Pydantic v2, Python 3.11/3.12+.",
        "Database & Storage: PostgreSQL 16 (w/ pgvector), MinIO (S3-compatible).",
        "Asynchronous Workers: Celery + Redis for heavy AI tasks.",
        "AI Stack: Whisper (faster-whisper), pyannote.audio, Ollama.",
    ],
)

# Slide 6: Current Progress (What's Done)
add_bullet_slide(
    "Current Progress (Review 1)",
    [
        "Finalized Product Requirements & Technical Architecture (Docs & Diagrams).",
        "Frontend CI/CD & Foundation: Next.js setup, linting, formatting resolved.",
        "Backend Core: FastAPI established, database models configured, CI pipelines passing.",
        "Streaming Pipeline: Implemented wss:// WebSocket endpoints for live audio ingestion.",
        "AI Stubbing: Created modular Protocol interfaces for STT, LLM, and Diarization services.",
        "Chrome Extension: Scaffolded background & offscreen capture infrastructure.",
    ],
)

# Slide 7: Next Steps & Roadmap
add_bullet_slide(
    "Next Steps",
    [
        "Infrastructure Integration: Finalize Docker Compose, MinIO, and Celery worker setups.",
        "AI Model Integration: Connect the modular stubs to actual local Whisper and Ollama instances.",
        "Frontend Integration: Build the dashboard to view transcripts and RAG capabilities.",
        "End-to-End Testing: Complete an end-to-end flow from Chrome Extension capture to AI Summary.",
    ],
)

# Save the presentation
prs.save("../../MeetingMind_Internal_Review_1.pptx")
print("Presentation generated successfully at MeetingMind_Internal_Review_1.pptx")
